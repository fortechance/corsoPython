"""SCRIPT

-CONVERTE FILE DI CONFIGURAZIONE DELLE CPE (.txt) IN FILE JINJA2 .j2
-LEGGE DATI ORDINE DA DELIVERY FORM E CONTROLLA ASSENZA/CONGRUITA' DATI
-CREA FILE DI CONFIGURAZIONE .conf SU BASE SERVIZIO CON DATI DELIVERY FORM + INTEGRAZIONI DERIVATE
-COMPILA FILE .xls
-COMPILA FILE .pptx

"""

#TODO PM: INSERIRE CAMPI VUOTI PER VLAN ID ALTRIMENTI USER INPUT
#TODO TIPO VCLOUD: STD, STD2, DIRECTOR???
#TODO ALTRI SERVIZI? VPDC, ETC.
#TODO RIMUOVERE FUNZIONE deliv_form_to_dict()?
#TODO CHECK FILE DI CONFIGURAZIONE .conf
#TODO CONFIGURAZIONI APPARATI VIA TELNET/SSH CON CHECK PRE E POST (pyats?)

#IMPORT
from os.path import join
import re
import glob
import os
import pandas as pd
pd.set_option('display.max_rows', None)
# from collections import defaultdict
from pprint import pprint
from netaddr import *
from jinja2 import Environment, FileSystemLoader, Template, select_autoescape
import warnings
warnings.simplefilter(action='ignore', category=FutureWarning)
from tabulate import tabulate

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
# from pptx.enum.shapes import MSO_SHAPE
from openpyxl import load_workbook

#VARIABILI GLOBALI
# dict_deliv_form = defaultdict(lambda: "KEY_NOT_DEFINED")
dict_deliv_form = {}
campi_list_of_dicts = []
dfintegrato = {}
dict_for_slides = {}
datamiss = []

# CREA VARIABILE CARTELLA DI QUESTO SCRIPT 
os.chdir(os.path.dirname(__file__))
scriptdir = os.getcwd()


# SI SPOSTA NELLA CARTELLA "Dati" E RECUPERA DELIVERY FORM XLSX
os.chdir(join(scriptdir, "Dati")) 
dirdati = os.getcwd()
delivery_form = glob.glob(join(dirdati, "*.xlsm"))
delivery_form = delivery_form[0].replace("\\\\", "\\")


# CREA DF PANDAS DA FILE EXCEL DELIVERY FORM
df = pd.read_excel(delivery_form, sheet_name="HandOver", usecols="C,D", names = ["CAMPO", "VALORE"], skiprows=2)
# df = df.set_index('CAMPO')
df = df.fillna(0)  # senza questo da errore la conversione sotto (TBV)

# CONVERTE I CAMPI FLOAT IN INT ALTRIMENTI VENGONO NUMERI DECIMALI ANCHE SE NON CI SONO TBV
for (columnName, columnData) in df.iteritems():
    # print(columnName, columnData.dtype)
    if columnData.dtype == "float64":
        df[columnName] = df[columnName].astype(int)


# SI "AGGIUSTANO" I CAMPI DEL DELIVERY FORM
df['VALORE'] = df['VALORE'].replace(0, "Empty!")
df = df.replace(0, "")
df.CAMPO = df.CAMPO.str.strip()
df.CAMPO = df.CAMPO.str.replace(" ","_") #le variabili nei template jinja2 sono senza spazi
df.CAMPO = df.CAMPO.str.replace("⌞", "") 
df.CAMPO = df.CAMPO.str.replace("(", "") 
df.CAMPO = df.CAMPO.str.replace(")", "") 
df.CAMPO = df.CAMPO.str.replace("[", "") 
df.CAMPO = df.CAMPO.str.replace("]", "")
df.CAMPO = df.CAMPO.str.replace("/", "")
df.CAMPO = df.CAMPO.str.replace("?", "")
df.CAMPO = df.CAMPO.str.replace(":", "")
df.CAMPO = df.CAMPO.str.replace("n°", "num")
df.CAMPO = df.CAMPO.str.replace("N°", "num")
df.CAMPO = df.CAMPO.str.replace("tà", "ta")
df.CAMPO = df.CAMPO.str.replace("_Quantitativo", "Quantitativo")
df.CAMPO = df.CAMPO.str.replace("_Tipologia", "Tipologia")
df.CAMPO = df.CAMPO.str.replace("È", "")
df.CAMPO = df.CAMPO.str.replace("è", "")
df.CAMPO = df.CAMPO.str.replace("E-_mail", "E_mail")
df.CAMPO = df.CAMPO.str.replace("_richiest", "richiest")
df.CAMPO = df.CAMPO.str.replace("_previsto", "previsto")
df.CAMPO = df.CAMPO.str.replace("n._IP_Pubblici", "num")
df.CAMPO = df.CAMPO.str.replace("sarà", "viene")
df.CAMPO = df.CAMPO.str.replace("Q.ta", "Qta")
df.CAMPO = df.CAMPO.str.replace("verrà", "viene")
df = df.astype(str)
#print(df.CAMPO.str.isalnum())
print(df)


#CREA TABELLA MAXCALLS
dfmaxcalls = pd.read_excel('maxcalls.xlsx', sheet_name="maxcalls")
# print(dfmaxcalls)

#CREA DICT DA DF
dftodict = df.set_index('CAMPO')
for entry in dftodict.columns:
        for tupla in dftodict[entry].items(): #.keys(), .values, .items()
            dict_deliv_form[tupla[0]] = tupla[1]
# print(dict_deliv_form)



#CONTROLLA PRESENZA VALORE CHIAVE RICHIESTA IN DICT_DELIV_FORM
def checklist_dati():
    
    #NOTE INSERIRE ANCHE LE VLAN SE LE INSERIRANNO NEL DELIVERY FORM
    check_vpdc_ucc = ['Ragione_Sociale_Cliente','Condivisa_o_Dedicata','Subnet_IP_pubblica_staccata_su_SOM','banda_Internet_Mbps','Subnet_Interconnessione','Prima_Subnet_IP_30_MPLS_Dati','Seconda_Subnet_IP_30_MPLS_Dati','Banda_MPLS_dati_Mbps','Subnet_IP_private_MPLS_Fonia','Prima_Subnet_IP_30_MPLS_Fonia','Seconda_Subnet_IP_30_MPLS_Fonia', 'Banda_MPLS_Fonia_Mbps','Tipologia_Servizio','Nome_VRF_Cliente_MPLS_Dati','RD_per_accessi_VPN_Dati','Nome_VRF_Cliente_MPLS_Fonia','RD_per_accessi_VPN_Fonia','LAN_SBC_remota','Simultaneus_Call', 'Tipologia_Firewall']
    
    check_vpdc_internet = []
    
    #ECCEZIONI

    
    for d in check_vpdc_ucc:
        if dict_deliv_form[d] == "Empty!":
            datamiss.append(d)

    if datamiss:
        print('\n\nATTENZIONE! FORNIRE I SEGUENTI DATI OBBLIGATORI PER PROCEDERE:')
        for i, e in enumerate(datamiss):
            print(i+1,') ', e)
        print('\n')
        
        print('KAM: ', dict_deliv_form['Referente_KAM_e-mail'])
        print('SCRIVERE EMAIL A: ', dict_deliv_form['Referente_PM_Nome__Cognome'])
        print(f"OGGETTO: Ordine {dict_deliv_form['Ragione_Sociale_Cliente']} (Codice SOM: {dict_deliv_form['Codice_SOM_UDP']})" )
        
        print(f"\nCiao { dict_deliv_form['Referente_PM_Nome__Cognome'].title()}, \n\npuoi fornirmi gentilmente i seguenti dati mancanti nel Delivery Form relativo a {dict_deliv_form['Ragione_Sociale_Cliente']} (Codice SOM: {dict_deliv_form['Codice_SOM_UDP']})?\n ")
        for i, e in enumerate(datamiss):
            print(i+1,') ', e)
        print('\nGrazie,\nSaluti.\n')
        quit()

#CONVERTE FILE .txt in .j2 PER CONFIGURAZIONI CPE/PE
def txt_template_to_j2():
    #CONVERTE FILE .txt in .j2
    #Il file di testo deve essere ".txt" e deve avere i campi tra "<" e ">")
    #Trova anche i campi dei file
    
    global campi_list_of_dicts
    
    #Si sposta nella cartella dei templates
    os.chdir(join(scriptdir, "Templates", "txt")) 
    txttemplatedir = os.getcwd()
    txttemplatedirfiles = os.listdir(txttemplatedir)
    #print(os.getcwd())
    
    #verifica se ci sono file di testo e li converte (sostituisce <campo> con {{campo}}, cambia estensione da .txt a .j2 e rimuove nel nome gli spazi)    
    for file in txttemplatedirfiles:

        if file.endswith(".txt"):
            
            #legge il file di testo e lo chiude
            with open(file, "r") as f:
                fcontent = f.read()
                #print(fcontent)

            #converte il contenuto del file di testo in file jinja2
            fcontent = fcontent.replace("<","{{").replace(">","}}")
            #print(fcontent)
            
            #scrive l'output in un nuovo file .j2 nella cartella "jinja"
            # os.chdir("../jinja/")
            j2filename = ''.join( file.rpartition(".")[:-2])
            
            #normalizza nome file
            j2filename = j2filename.replace(" ","_").replace("_-_","_").replace("-_","_").replace(".","_")
            # print(j2filename)
            
            #path file
            j2filenamepath = join(scriptdir, "Templates", "jinja", j2filename) + ".j2"
            # print(j2filenamepath)
            
            with open(j2filenamepath, "w") as f:
                f.write(fcontent)
        
            #cerca i campi presenti nei file e crea dict (campi duplicati rimossi)
            campi = re.findall("\{\{(\w+)\}\}", fcontent)
            # parent_dict = j2filename
            
            # crea lista di dizionari con i campi per ciascun file
            dict_from_key = dict.fromkeys(campi) #campi nel file
            dict_from_key['filename'] = j2filename #aggiunge key "filename" con nome file.j2
            campi_list_of_dicts.append(dict_from_key) #aggiunge il dizionario alla lista di dizionari
            #campi_dict_of_dicts[j2filename + ".j2"] = dict.fromkeys(campi) #alternativa dict of dicts
            # print(campi_list_of_dicts)
            
        else:
            if os.path.isdir(file) == False: 
                #riporta eventuale estensione non corretta
                # print(f'ATTENZIONE, ESTENSIONE "{file.rsplit(".",1)}" NON CORRETTA. NOME FILE:', file)
                estensione = ''.join(file.rpartition(".")[-2:]) #join della tupla di output (punto, estensione)
                print(f'ATTENZIONE, ESTENSIONE "{estensione}" NON CORRETTA. MODIFICARE IL NOME DEL FILE:', file)




#INTEGRA IL dict DEL DELIVERY FORM
def deliv_form_integrazione():
    
    global dfintegrato
    global dict_for_slides
    global datamiss
    
    #lista datamiss contiene i dati mancanti fondamentali per la creazione delle configurazioni

        
    #DATI COMUNI
    nome_cliente = dict_deliv_form['Ragione_Sociale_Cliente']
    
    #nome client trim
    nome_cliente_ = nome_cliente.replace('--', '-').replace('  ', ' ').replace('__', '_').replace('..','.').replace('//','/').replace(' ','_').strip().upper()
    dict_deliv_form['Nome_Cliente'] = nome_cliente_
    
    #NOTE VCLOUD_UC_STD O VCLOUD_DIRECTOR_CLIENTI??? USER INPUT?
    vcloud = 'VCLOUD_UC_STD-2' 
    
    tenant = vcloud + '|' + "APP_" + nome_cliente_
    dict_deliv_form['Tenant'] = tenant
    # print(tenant)
    
    
    ###############################
    ###FILE1: ACI-VPDC Internet.txt
    ###############################
    
    #NOTE VLAN RICHIEDE USER INPUT
    dict_deliv_form['Vlan_internet'] = '1041'
    
    #INDIRIZZI IP
    subnet_IP_pubblica_staccata_su_SOM = IPNetwork(dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM']) # network/bitmask
    
    #INTEGRA NEL DICT INDIRIZZI IP 
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_NETWORK'] = str(subnet_IP_pubblica_staccata_su_SOM.network) # network
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_HSRP'] = str(subnet_IP_pubblica_staccata_su_SOM[1])
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_PRI'] = str(subnet_IP_pubblica_staccata_su_SOM[2])
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_SEC'] = str(subnet_IP_pubblica_staccata_su_SOM[3])
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_NSXEDGE'] = str(subnet_IP_pubblica_staccata_su_SOM[4])
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_NETMASK'] = str(subnet_IP_pubblica_staccata_su_SOM.netmask)
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_BITMASK'] = str(subnet_IP_pubblica_staccata_su_SOM.prefixlen)
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_BROADCAST'] = str(subnet_IP_pubblica_staccata_su_SOM.broadcast)
    dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_WILDMASK'] = str(subnet_IP_pubblica_staccata_su_SOM.hostmask)
    
    #INTEGRA NEL DICT ALTRE INFO
    dict_deliv_form['epg_internet'] = 'APP_' + nome_cliente_ + "|" + "EPG_INTERNET"
    dict_deliv_form['vcloud_epg_internet'] = vcloud + "|" + nome_cliente_ + "|" + "EPG_INTERNET"
    dict_deliv_form['banda_Internet_Mbps'] = dict_deliv_form['banda_Internet_Mbps'].replace(' Mbps','')
    
    #NOTE SEQUENCE NUMBER PREFIX-LIST CONNECTED (E STATIC?) RICHIEDE USER INPUT
    dict_deliv_form['ip_prefix_list_connected_seq'] = 'seq_connected_user_input'
    dict_deliv_form['ip_prefix_list_static_seq'] = 'seq_static_user_input'
    
    #PER SLIDE/XLSX
    dict_deliv_form['epg_internet_slide'] = 'APP_' + nome_cliente_ + "|" + "EPG_INTERNET"    
    dict_deliv_form['epg_internet_xlsx'] = vcloud + '|APP_' + nome_cliente_ + "|" + "EPG_INTERNET"    
    
    
    #######################################################
    ###FILE2: Cliente - ACI-VPDC FONIA_mpls_asr_qos_QoS.txt
    #######################################################
    
    #NOTE VLAN RICHIEDE USER INPUT
    dict_deliv_form['Vlan_Subnet_Interconnessione'] = '1051'
    
    #INDIRIZZI IP
    subnet_Interconnessione = IPNetwork(dict_deliv_form['Subnet_Interconnessione']) # network/bitmask
    prima_Subnet_IP_30_MPLS_Dati = IPNetwork(dict_deliv_form['Prima_Subnet_IP_30_MPLS_Dati']) # network/bitmask
    seconda_Subnet_IP_30_MPLS_Dati = IPNetwork(dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Dati']) # network/bitmask
    
    #INTEGRA NEL DICT INDIRIZZI IP 
    dict_deliv_form['Subnet_Interconnessione_NETWORK'] = str(subnet_Interconnessione.network) # network
    dict_deliv_form['Prima_Subnet_IP_30_MPLS_Dati_NETWORK'] = str(prima_Subnet_IP_30_MPLS_Dati.network) # network
    dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Dati_NETWORK'] = str(seconda_Subnet_IP_30_MPLS_Dati.network) # network
    dict_deliv_form['Subnet_Interconnessione_HSRP'] = str(subnet_Interconnessione[1])
    dict_deliv_form['Subnet_Interconnessione_PRI'] = str(subnet_Interconnessione[2])
    dict_deliv_form['Subnet_Interconnessione_SEC'] = str(subnet_Interconnessione[3])
    dict_deliv_form['Subnet_Interconnessione_vFW_UCC'] = str(subnet_Interconnessione[-2])
    dict_deliv_form['Subnet_Interconnessione_NETMASK'] = str(subnet_Interconnessione.netmask)
    dict_deliv_form['Subnet_Interconnessione_BITMASK'] = str(subnet_Interconnessione.prefixlen)
    dict_deliv_form['Subnet_Interconnessione_BROADCAST'] = str(subnet_Interconnessione.broadcast)
    dict_deliv_form['Subnet_Interconnessione_WILDMASK'] = str(subnet_Interconnessione.hostmask)
    dict_deliv_form['Prima_Subnet_IP_30_MPLS_Dati_CPE_PRI'] =	str(prima_Subnet_IP_30_MPLS_Dati[1])	
    dict_deliv_form['Prima_Subnet_IP_30_MPLS_Dati_PE_PRI']	=	str(prima_Subnet_IP_30_MPLS_Dati[2])
    dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Dati_CPE_SEC'] = str(seconda_Subnet_IP_30_MPLS_Dati[1])
    dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Dati_PE_SEC'] = str(seconda_Subnet_IP_30_MPLS_Dati[2])
	
    #INTEGRA NEL DICT ALTRE INFO
    dict_deliv_form['Banda_MPLS_dati_Mbps'] = dict_deliv_form['Banda_MPLS_dati_Mbps'].replace(' Mbps','')
    
    #RITORNA LA BANDA VOIP MPLS DALLA TABELLA MAX CALLS
    banda_MPLS = int(dict_deliv_form['Banda_MPLS_dati_Mbps'])
    
    #Ci sono due "BANDA MPLS a 2Mbps" ma la BW VoIP MPLS è sempre 1024 per cui si sceglie il primo risultato .tolist[0]
    if banda_MPLS == 2: 
        banda_VoIP_MPLS_Garantita = dfmaxcalls.loc[dfmaxcalls['Banda_MPLS' ] == banda_MPLS, 'Banda_VoIP_MPLS_Garantita'].tolist()[0]
        dict_deliv_form['Banda_VoIP_MPLS_Garantita'] = banda_VoIP_MPLS_Garantita
    else: #qui è un .item()
        banda_VoIP_MPLS_Garantita = dfmaxcalls.loc[dfmaxcalls['Banda_MPLS' ] == banda_MPLS, 'Banda_VoIP_MPLS_Garantita'].item()
        dict_deliv_form['Banda_VoIP_MPLS_Garantita'] = banda_VoIP_MPLS_Garantita

    dict_deliv_form['MAXIMUM_ROUTES_MPLS_Dati'] = '500'
    
    #PER CONF/SLIDE/XLS
    dict_deliv_form['epg_vpn_mpls_conf'] = 'APP_' + nome_cliente_ + "|" + "EPG_UC"
    dict_deliv_form['epg_vpn_mpls_slide'] = dict_deliv_form['epg_vpn_mpls_conf']
    dict_deliv_form['epg_vpn_mpls_xlsx'] = vcloud + '|APP_' + nome_cliente_ + "|" + "EPG_UC"
    
    
    ##############################################
    ###FILE3: Cliente - ACI-VPDC FONIA L3 vrf voce
    ##############################################
    
    #NOTE VLAN RICHIEDE USER INPUT
    dict_deliv_form['Vlan_IP_private_MPLS_Fonia'] = '1040'
    
    #INDIRIZZI IP
    subnet_ip_private_mpls_fonia = IPNetwork(dict_deliv_form['Subnet_IP_private_MPLS_Fonia']) # network/bitmask
    prima_Subnet_IP_30_MPLS_Fonia = IPNetwork(dict_deliv_form['Prima_Subnet_IP_30_MPLS_Fonia']) # network/bitmask
    seconda_Subnet_IP_30_MPLS_Fonia = IPNetwork(dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Fonia']) # network/bitmask
    
    #INTEGRA NEL DICT INDIRIZZI IP 
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_NETWORK'] = str(subnet_ip_private_mpls_fonia.network) # network
    dict_deliv_form['Prima_Subnet_IP_30_MPLS_Fonia_NETWORK'] = str(prima_Subnet_IP_30_MPLS_Fonia.network) # network
    dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Fonia_NETWORK'] = str(seconda_Subnet_IP_30_MPLS_Fonia.network) # network
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_HSRP'] = str(subnet_ip_private_mpls_fonia[1])
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_PRI'] = str(subnet_ip_private_mpls_fonia[2])
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_SEC'] = str(subnet_ip_private_mpls_fonia[3])
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_NSXEDGE'] = str(subnet_ip_private_mpls_fonia[10])
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_NETMASK'] = str(subnet_ip_private_mpls_fonia.netmask)
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_BITMASK'] = str(subnet_ip_private_mpls_fonia.prefixlen)
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_BROADCAST'] = str(subnet_ip_private_mpls_fonia.broadcast)
    dict_deliv_form['Subnet_IP_private_MPLS_Fonia_WILDMASK'] = str(subnet_ip_private_mpls_fonia.hostmask)
    dict_deliv_form['Prima_Subnet_IP_30_MPLS_Fonia_CPE_PRI'] =	str(prima_Subnet_IP_30_MPLS_Fonia[1])	
    dict_deliv_form['Prima_Subnet_IP_30_MPLS_Fonia_PE_PRI']	=	str(prima_Subnet_IP_30_MPLS_Fonia[2])
    dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Fonia_CPE_SEC'] = str(seconda_Subnet_IP_30_MPLS_Fonia[1])
    dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Fonia_PE_SEC'] = str(seconda_Subnet_IP_30_MPLS_Fonia[2])
                                                            
    #INTEGRA NEL DICT ALTRE INFO
    dict_deliv_form['Banda_MPLS_Fonia_Mbps'] = dict_deliv_form['Banda_MPLS_Fonia_Mbps'].replace(' Mbps','')
    
    #RITORNA LA BANDA VOIP ITCFONIA DALLA TABELLA MAX CALLS
    banda_MPLS_Fonia_Mbps = int(dict_deliv_form['Banda_MPLS_Fonia_Mbps'])
    banda_VoIP_ITCFONIA_Garantita = dfmaxcalls.loc[dfmaxcalls['Banda_ITCFONIA'] == banda_MPLS_Fonia_Mbps, 'Banda_VoIP_ITCFONIA_Garantita'].item() #cerca associazione nella tabella maxcalls
    dict_deliv_form['Banda_VoIP_ITCFONIA_Garantita'] = banda_VoIP_ITCFONIA_Garantita
    
    dict_deliv_form['MAXIMUM_ROUTES'] = '500'
    
    #PER SLIDE
    dict_deliv_form['epg_ictfonia_slide'] = 'APP_' + nome_cliente_ + "|" + "EPG_ITCFONIA"
    
    #PRINT DICT
    #pprint(dict_deliv_form,sort_dicts=True)


    #RICREA df CON DATI COMPLETI E STAMPA IN FILE .txt o .xls I
    dfintegrato = pd.DataFrame.from_dict(dict_deliv_form,orient='index')
    # print(tabulate(dfintegrato,tablefmt='grid'))
    
    #file di testo
    filedatitxt = join(scriptdir, "03_DOC_NETWORK", dict_deliv_form['Ragione_Sociale_Cliente'] + ' - CS_' + dict_deliv_form['Codice_SOM_UDP'] + '_ID_' + dict_deliv_form['ID_Ordine_Commerciale'] + '.txt')
    
    with open(filedatitxt, "w") as f:
        f.write(tabulate(dfintegrato,tablefmt='grid')) 
        # f.write(tabulate(dfintegrato,tablefmt='html')) 
        # for k,v in dict_deliv_form.items():
        #     f.write(f"{k}: {v}\n")
        # f.write(tabulate(df.to_string(header=False))) #max_colwidth=50
    
    #file excel
    filedatixls =  join(scriptdir, "03_DOC_NETWORK", dict_deliv_form['Ragione_Sociale_Cliente'] + ' - CS_' + dict_deliv_form['Codice_SOM_UDP'] + '_ID_' + dict_deliv_form['ID_Ordine_Commerciale'] + '.xlsx')
           
    dfintegrato.to_excel(filedatixls, sheet_name="ORDINE", merge_cells=False)  # crea l'excel
        

#CREA I FILE DI CONFIGURAZIONE CON I JINJA2 TEMPLATES
def jinja_config():
    
    #NOTE PREPARA FILE SU BASE SERVIZIO "Tipologia_Servizio"
    #print(os.getcwd())

    #ENVIROMENT JINJIA2. INDICARE LA CARTELLA IN CUI SONO CONTENUTI I TEMPLATE
    #https://jinja.palletsprojects.com/en/3.0.x/api/#jinja2.Environment
    env = Environment(
    loader=FileSystemLoader(
            join(scriptdir, "Templates", "jinja")
    ),
    keep_trailing_newline=True, #True = mantiene un eventuale riga finale vuota presente nel template. False la rimuove
    autoescape=select_autoescape())
    
    #SERVIZI (Tipologia Servizio)
    # vPDC_UCC = ['Cliente_ACI-VPDC_FONIA_L3_vrf_voce.j2','Cliente_ACI-VPDC_Internet.j2','Cliente_ACI-VPDC_FONIA_mpls_asr_qos_QoS.j2']
    tipologia_servizio = dict_deliv_form['Tipologia_Servizio']
    
    #SELEZIONE FILE DI CONFIG SU BASE SERVIZIO DELIVERY FORM
    if tipologia_servizio == "vPDC UCC":
        
        config_file_servizio = ['Cliente_ACI-VPDC_FONIA_L3_vrf_voce_OK.j2','Cliente_ACI-VPDC_Internet_OK.j2', 'Cliente_ACI-VPDC_FONIA_mpls_asr_qos_QoS_OK.j2']
        
    elif tipologia_servizio == "vPDC":
        
        config_file_servizio = []
    
    else:
        print('SERVIZIO SCONOSCIUTO...ESCO!')
        quit()
    
    #LOOP RENDERING CONFIGURAZIONE IN jinja2
    for f in config_file_servizio:
        
        #Seleziona un file template .j2 dalla cartella di cui sopra
        template = env.get_template(f)
    
        # conf = join("../03_DOC_NETWORK","Cliente_ACI-VPDC_FONIA_L3_vrf_voce.conf")
        configfile = join(scriptdir, "03_DOC_NETWORK", "NETWORK", f.replace('Cliente_', dict_deliv_form['Ragione_Sociale_Cliente'] + ' - ').replace("j2","conf"))
    
        with open(configfile, "w") as f:
            f.write(template.render(dict_deliv_form))


#COMPILA TEMPLATE PPTX
def compila_pptx():
    
    #COMPILA TEMPLATE PPTX
    os.chdir(join(scriptdir, "Templates", "pptx"))
    
    vpdc = Presentation('VPDC UCC.pptx')
    
    #TITOLO SLIDE INIZIALE
    for shape in vpdc.slides[0].shapes:
        if 'RAGIONE SOCIALE' in shape.text:
            shape.text_frame.paragraphs[0].text = dict_deliv_form['Ragione_Sociale_Cliente']
            shape.text_frame.paragraphs[0].font.name = 'Lucida Grande'
            shape.text_frame.paragraphs[0].font.size = Pt(64)

    #TITOLO SLIDE 1
    for shape in vpdc.slides[1].shapes:
        if shape.has_text_frame:
            if 'RAGIONE SOCIALE' in shape.text:
                shape.text_frame.paragraphs[0].text = dict_deliv_form['Ragione_Sociale_Cliente'] + ' - ACI VPDC UCC vCloud'
                shape.text_frame.paragraphs[0].font.name = 'Lucida Grande'
                shape.text_frame.paragraphs[0].font.size = Pt(28)
    
    #TITOLO SLIDE 3
    #NOTE TBC
    
    
    #INDIVIDUA LE TABELLE DELLA SLIDE1
    slide = vpdc.slides[1]
    
    for shape in slide.shapes: #for tra le shape della slide
        
        if shape.has_table: #se la forma è una tabella
            
            #vars celle tabelle (cosmetico)
            c_0_0 = 'shape.table.cell(0,0).text_frame.paragraphs[0]'
            c_1_0 = 'shape.table.cell(1,0).text_frame.paragraphs[0]'
            c_1_1 = 'shape.table.cell(1,1).text_frame.paragraphs[0]'
            c_2_0 = 'shape.table.cell(2,0).text_frame.paragraphs[0]'
            c_2_1 = 'shape.table.cell(2,1).text_frame.paragraphs[0]'
            c_3_0 = 'shape.table.cell(3,0).text_frame.paragraphs[0]'
            c_3_1 = 'shape.table.cell(3,1).text_frame.paragraphs[0]'
            c_4_0 = 'shape.table.cell(4,0).text_frame.paragraphs[0]'
            c_4_1 = 'shape.table.cell(4,1).text_frame.paragraphs[0]'
            c_4_2 = 'shape.table.cell(4,2).text_frame.paragraphs[0]'
        

            for i, cella in enumerate(shape.table.iter_cells()): #for tra le celle della tabella
                # print(i, ": ", cella.text) #stampa il testo di tutte le celle della tabella
                
                #NOTE LA FORMATTAZIONE SI APPLICA A LIVELLO DI PARAGRAFO O DI RUN (singola lettera) (VEDI LINK INIZIO SCRIPT)
                
                #COLORI
                red = RGBColor(255, 0, 0)
                marrone = RGBColor(134, 94, 61)
                
                
                #TENANT
                if 'TENANT' in eval(c_0_0).text:
                    
                    #INSERIMENTO DATI TENANT
                    eval(c_1_0).text = dict_deliv_form['Tenant']
                    
                    #FORMATTAZIONE TENANT
                    epg_format_tenant = {'font_name': 'Arial', 'font_italic': False, 'font_bold': True, 'font_size': Pt(10), 'font_underline' : False, 'font_color': red , 'allineamento': PP_ALIGN.LEFT}
                    
                    eval(c_1_0).font.name = epg_format_tenant['font_name']
                    eval(c_1_0).font.italic = epg_format_tenant['font_italic']
                    eval(c_1_0).font.bold = epg_format_tenant['font_bold'] 
                    eval(c_1_0).font.size = epg_format_tenant['font_size'] 
                    eval(c_1_0).font.underline = epg_format_tenant['font_underline']
                    eval(c_1_0).font.color.rgb = epg_format_tenant['font_color']
                
                    
                #EPGs
                    #INSERIMENTO DATI EPG
                if  'EPG_INTERNET' in shape.table.cell(0,0).text:
                    shape.table.cell(0,0).text = dict_deliv_form['epg_internet']
                    shape.table.cell(1,0).text = dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM']
                    shape.table.cell(1,1).text = "vlan-" + dict_deliv_form['Vlan_internet']
                    
                if  'EPG_UC' in shape.table.cell(0,0).text:
                    shape.table.cell(0,0).text = dict_deliv_form['epg_vpn_mpls_slide']
                    shape.table.cell(1,0).text = dict_deliv_form['Subnet_Interconnessione']
                    shape.table.cell(1,1).text = "vlan-" + dict_deliv_form['Vlan_Subnet_Interconnessione']
                
                if  'EPG_ITCFONIA' in shape.table.cell(0,0).text:
                    shape.table.cell(0,0).text = dict_deliv_form['epg_ictfonia_slide']
                    shape.table.cell(1,0).text = dict_deliv_form['Subnet_IP_private_MPLS_Fonia']
                    shape.table.cell(1,1).text = "vlan-" + dict_deliv_form['Vlan_IP_private_MPLS_Fonia']
                
                    #FORMATTAZIONE EPG
                if 'EPG_' in shape.table.cell(0,0).text:
                    
                    epg_format_c_0_0 = {'font_name': 'Calibri', 'font_italic': False, 'font_bold': False, 'font_size': Pt(14), 'font_underline' : False}
                    epg_format_c_1_0 = {'font_name': 'Calibri', 'font_italic': False, 'font_bold': True, 'font_size': Pt(14), 'font_underline' : False}
                    epg_format_c_1_1 = {'font_name': 'Calibri', 'font_italic': False, 'font_bold': True, 'font_size': Pt(14), 'font_underline' : False}
                    
                    eval(c_0_0).font.name = epg_format_c_0_0['font_name']
                    eval(c_0_0).font.italic = epg_format_c_0_0['font_italic']
                    eval(c_0_0).font.bold = epg_format_c_0_0['font_bold'] 
                    eval(c_0_0).font.size = epg_format_c_0_0['font_size'] 
                    eval(c_0_0).font.underline = epg_format_c_0_0['font_underline']
                    # eval(c_1_0).font.color = ''
                    #
                    eval(c_1_0).font.name = epg_format_c_1_0['font_name']
                    eval(c_1_0).font.italic = epg_format_c_1_0['font_italic']
                    eval(c_1_0).font.bold = epg_format_c_1_0['font_bold'] 
                    eval(c_1_0).font.size = epg_format_c_1_0['font_size'] 
                    eval(c_1_0).font.underline = epg_format_c_1_0['font_underline']
                    # eval(c_1_0).font.color = ''
                    #
                    eval(c_1_1).font.name = epg_format_c_1_1['font_name']
                    eval(c_1_1).font.italic = epg_format_c_1_1['font_italic']
                    eval(c_1_1).font.bold = epg_format_c_1_1['font_bold'] 
                    eval(c_1_1).font.size = epg_format_c_1_1['font_size'] 
                    eval(c_1_1).font.underline = epg_format_c_1_1['font_underline']
                    # eval(c_1_0).font.color = ''
                
                #VPN (INTERNET, MPLS, ITCFONIA)
                if  'VPN INTERNET' in shape.table.cell(0,0).text:
                    
                    #INSERIMENTO DATI VPN INTERNET
                    shape.table.cell(1,0).text = 'Banda' + ' ' + dict_deliv_form['banda_Internet_Mbps'] + ' Mbps'
                    
                    #FORMATTAZIONE VPN INTERNET
                    epg_format_vpn_int_c_0_0 = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': True, 'font_size': Pt(16), 'font_underline' : False, 'allineamento': PP_ALIGN.CENTER}
                    epg_format_vpn_int_c_1_0 = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': True, 'font_size': Pt(10), 'font_underline' : True, 'allineamento': PP_ALIGN.CENTER}
                    #
                    eval(c_0_0).font.name = epg_format_vpn_int_c_0_0['font_name']
                    eval(c_0_0).font.italic = epg_format_vpn_int_c_0_0['font_italic']
                    eval(c_0_0).font.bold = epg_format_vpn_int_c_0_0['font_bold'] 
                    eval(c_0_0).font.size = epg_format_vpn_int_c_0_0['font_size'] 
                    eval(c_0_0).font.underline = epg_format_vpn_int_c_0_0['font_underline']
                    eval(c_0_0).alignment = epg_format_vpn_int_c_0_0['allineamento']
                    #
                    eval(c_1_0).font.name = epg_format_vpn_int_c_1_0['font_name']
                    eval(c_1_0).font.bold = epg_format_vpn_int_c_1_0['font_bold'] 
                    eval(c_1_0).font.size = epg_format_vpn_int_c_1_0['font_size']
                    eval(c_1_0).font.underline = epg_format_vpn_int_c_1_0['font_underline']
                    eval(c_1_0).alignment = epg_format_vpn_int_c_1_0['allineamento']

                    #INSERIMENTO DATI VPN MPLS
                if  'VPN MPLS' in shape.table.cell(0,0).text:

                    shape.table.cell(1,1).text = dict_deliv_form['Nome_VRF_Cliente_MPLS_Dati']
                    shape.table.cell(2,1).text = dict_deliv_form['RD_per_accessi_VPN_Dati']
                    shape.table.cell(3,0).text = 'Banda ' + dict_deliv_form['Banda_MPLS_dati_Mbps'] + ' Mbps'
                    shape.table.cell(4,0).text = 'Banda fonia garantita'
                    shape.table.cell(4,2).text = str(dict_deliv_form['Banda_VoIP_MPLS_Garantita']) + ' Kbps'

                    #INSERIMENTO DATI VPN ITCFONIA
                if  'VPN ITCFONIA' in shape.table.cell(0,0).text:
                    
                    shape.table.cell(1,1).text = dict_deliv_form['Nome_VRF_Cliente_MPLS_Fonia']
                    shape.table.cell(2,1).text = dict_deliv_form['RD_per_accessi_VPN_Fonia']
                    shape.table.cell(3,0).text = 'Banda ' + dict_deliv_form['Banda_MPLS_Fonia_Mbps'] + ' Mbps'
                    shape.table.cell(4,0).text = 'Banda fonia garantita'
                    shape.table.cell(4,2).text = str(dict_deliv_form['Banda_VoIP_ITCFONIA_Garantita']) + ' Kbps'
                
                    #FORMATTAZIONE VPN MPLS E VPN ITCFONIA
                if  'VPN MPLS' in shape.table.cell(0,0).text or 'VPN ITCFONIA' in shape.table.cell(0,0).text:
                    
                    epg_format_vpns_c_0_0 = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': True, 'font_size': Pt(16), 'font_underline' : False, 'allineamento': PP_ALIGN.CENTER}
                    epg_format_vpns_c_12_0 = {'font_name': 'Lucida Grande', 'font_italic': True, 'font_bold': False, 'font_size': Pt(10), 'font_underline' : False, 'allineamento': PP_ALIGN.LEFT}
                    epg_format_vpns_c_12_1 = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': True, 'font_size': Pt(10), 'font_underline' : False, 'allineamento': PP_ALIGN.LEFT}
                    epg_format_vpns_c_3_0 = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': True, 'font_size': Pt(10), 'font_underline' : True, 'allineamento': PP_ALIGN.CENTER}
                    epg_format_vpns_c_4_0 = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': False, 'font_size': Pt(10), 'font_underline' : False,  'allineamento': PP_ALIGN.LEFT}
                    epg_format_vpns_c_4_2 = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': True, 'font_size': Pt(11), 'font_underline' : False, 'font_color': marrone , 'allineamento': PP_ALIGN.LEFT}
                    #
                    eval(c_0_0).font.name = epg_format_vpns_c_0_0['font_name']
                    eval(c_0_0).font.italic = epg_format_vpns_c_0_0['font_italic']
                    eval(c_0_0).font.bold = epg_format_vpns_c_0_0['font_bold'] 
                    eval(c_0_0).font.size = epg_format_vpns_c_0_0['font_size'] 
                    eval(c_0_0).font.underline = epg_format_vpns_c_0_0['font_underline']
                    eval(c_0_0).alignment = epg_format_vpns_c_0_0['allineamento']
                    #
                    eval(c_1_0).font.name = eval(c_2_0).font.name = epg_format_vpns_c_12_0['font_name']
                    eval(c_1_0).font.bold = eval(c_2_0).font.bold = epg_format_vpns_c_12_0['font_bold']
                    eval(c_1_0).font.italic = eval(c_2_0).font.italic = epg_format_vpns_c_12_0['font_italic']
                    eval(c_1_0).font.size = eval(c_2_0).font.size = epg_format_vpns_c_12_0['font_size']
                    eval(c_1_0).font.underline = eval(c_2_0).font.underline = epg_format_vpns_c_12_0['font_underline']
                    eval(c_1_0).alignment = eval(c_2_0).alignment = epg_format_vpns_c_12_0['allineamento']
                    #
                    eval(c_1_1).font.name = eval(c_2_1).font.name = epg_format_vpns_c_12_1['font_name']
                    eval(c_1_1).font.bold = eval(c_2_1).font.bold = epg_format_vpns_c_12_1['font_bold']
                    eval(c_1_1).font.italic = eval(c_2_1).font.italic = epg_format_vpns_c_12_1['font_italic']
                    eval(c_1_1).font.size = eval(c_2_1).font.size = epg_format_vpns_c_12_1['font_size']
                    eval(c_1_1).font.underline = eval(c_2_1).font.underline = epg_format_vpns_c_12_1['font_underline']
                    eval(c_1_1).alignment = eval(c_2_1).alignment = epg_format_vpns_c_12_1['allineamento']
                    #
                    eval(c_3_0).font.name = epg_format_vpns_c_3_0['font_name']
                    eval(c_3_0).font.bold = epg_format_vpns_c_3_0['font_bold']
                    eval(c_3_0).font.italic = epg_format_vpns_c_3_0['font_italic']
                    eval(c_3_0).font.size = epg_format_vpns_c_3_0['font_size']
                    eval(c_3_0).font.underline = epg_format_vpns_c_3_0['font_underline']
                    eval(c_3_0).alignment = epg_format_vpns_c_3_0['allineamento']
                    #
                    eval(c_4_0).font.name = epg_format_vpns_c_4_0['font_name']
                    eval(c_4_0).font.bold = epg_format_vpns_c_4_0['font_bold']
                    eval(c_4_0).font.italic = epg_format_vpns_c_4_0['font_italic']
                    eval(c_4_0).font.size = epg_format_vpns_c_4_0['font_size']
                    eval(c_4_0).font.underline = epg_format_vpns_c_4_0['font_underline']
                    eval(c_4_0).alignment = epg_format_vpns_c_4_0['allineamento']
                    #
                    eval(c_4_2).font.name = epg_format_vpns_c_4_2['font_name']
                    eval(c_4_2).font.bold = epg_format_vpns_c_4_2['font_bold']
                    eval(c_4_2).font.italic = epg_format_vpns_c_4_2['font_italic']
                    eval(c_4_2).font.size = epg_format_vpns_c_4_2['font_size']
                    eval(c_4_2).font.underline = epg_format_vpns_c_4_2['font_underline']
                    eval(c_4_2).font.color.rgb = epg_format_vpns_c_4_2['font_color']
                    eval(c_4_2).alignment = epg_format_vpns_c_4_2['allineamento']
                    
                    
                    
                
                #HSRP/PRI/SEC
                    #INSERIMENTO DATI HSRP/PRI/SEC
                if  'HSRP internet' in shape.table.cell(0,0).text:
                    shape.table.cell(0,0).text = 'HSRP .' + dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_HSRP'].rpartition(".")[-1]
                    shape.table.cell(1,0).text = 'PRI .' + dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_PRI'].rpartition(".")[-1]
                    shape.table.cell(2,0).text = 'SEC .' + dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_SEC'].rpartition(".")[-1]
                    
                if  'HSRP mpls' in shape.table.cell(0,0).text:
                    shape.table.cell(0,0).text = 'HSRP .' + dict_deliv_form['Subnet_Interconnessione_HSRP'].rpartition(".")[-1]
                    shape.table.cell(1,0).text = 'PRI .' + dict_deliv_form['Subnet_Interconnessione_PRI'].rpartition(".")[-1]
                    shape.table.cell(2,0).text = 'SEC .' + dict_deliv_form['Subnet_Interconnessione_SEC'].rpartition(".")[-1]
                    
                if  'HSRP itcfonia' in shape.table.cell(0,0).text:
                    shape.table.cell(0,0).text = 'HSRP .' + dict_deliv_form['Subnet_IP_private_MPLS_Fonia_HSRP'].rpartition(".")[-1]
                    shape.table.cell(1,0).text = 'PRI .' + dict_deliv_form['Subnet_IP_private_MPLS_Fonia_PRI'].rpartition(".")[-1]
                    shape.table.cell(2,0).text = 'SEC .' + dict_deliv_form['Subnet_IP_private_MPLS_Fonia_SEC'].rpartition(".")[-1]
                
                    #FORMATTAZIONE CELLE HSRP/PRI/SEC
                if 'HSRP ' in shape.table.cell(0,0).text:
                    
                    epg_format_cpe = {'font_name': 'Lucida Grande', 'font_italic': False, 'font_bold': False, 'font_size': Pt(10), 'font_underline' : False}
                    
                    eval(c_0_0).font.name = eval(c_1_0).font.name = eval(c_2_0).font.name = epg_format_cpe['font_name']
                    eval(c_0_0).font.italic = eval(c_1_0).font.italic = eval(c_2_0).font.italic = epg_format_cpe['font_italic']
                    eval(c_0_0).font.bold = eval(c_1_0).font.bold = eval(c_2_0).font.bold = epg_format_cpe['font_bold'] 
                    eval(c_0_0).font.size = eval(c_1_0).font.size = eval(c_2_0).font.size = epg_format_cpe['font_size'] 
                    eval(c_0_0).font.underline = eval(c_1_0).font.underline = eval(c_2_0).font.underline = epg_format_cpe['font_underline']
                    # eval(c_1_0).font.color = ''
                

                #PUNTO-PUNTO
                
                    #INSERIMENTO DATI PtP            
                if  'ptp mpls' in shape.table.cell(0,0).text:
                    shape.table.cell(1,0).text = dict_deliv_form['Prima_Subnet_IP_30_MPLS_Dati']
                    shape.table.cell(2,0).text = dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Dati']
                
                if  'ptp itcfonia' in shape.table.cell(0,0).text:
                    shape.table.cell(1,0).text = dict_deliv_form['Prima_Subnet_IP_30_MPLS_Fonia']
                    shape.table.cell(2,0).text = dict_deliv_form['Seconda_Subnet_IP_30_MPLS_Fonia']
                
                    #FORMATTAZIONE PtP
                if 'ptp ' in shape.table.cell(0,0).text:
                    
                    epg_format_ptp = {'font_name': 'Calibri (Body)', 'font_italic': False, 'font_bold': False, 'font_size': Pt(12), 'font_underline' : False}
    
                    eval(c_0_0).font.name = eval(c_1_0).font.name = eval(c_2_0).font.name = epg_format_ptp['font_name']
                    eval(c_0_0).font.italic = eval(c_1_0).font.italic = eval(c_2_0).font.italic = epg_format_ptp['font_italic']
                    eval(c_0_0).font.bold = eval(c_1_0).font.bold = eval(c_2_0).font.bold = epg_format_ptp['font_bold'] 
                    eval(c_0_0).font.size = eval(c_1_0).font.size = eval(c_2_0).font.size = epg_format_ptp['font_size'] 
                    eval(c_0_0).font.underline = eval(c_1_0).font.underline = eval(c_2_0).font.underline = epg_format_ptp['font_underline']
                    # eval(c_1_0).font.color = ''
                    
            
                #SBC
                    #INSERIMENTO DATI SBC
                if  'SBC' in shape.table.cell(0,0).text:
                    shape.table.cell(0,0).text = 'SBC ' + dict_deliv_form['LAN_SBC_remota']
                    shape.table.cell(1,0).text = dict_deliv_form['Simultaneus_Call'] + ' contemporaneità'
                    
                    #FORMATTAZIONE SBC
                    epg_format_sbc = {'font_name': 'Calibri', 'font_italic': False, 'font_bold': False, 'font_size': Pt(14), 'font_underline' : False, 'allineamento': PP_ALIGN.CENTER}
    
                    eval(c_0_0).font.name = eval(c_1_0).font.name = epg_format_sbc['font_name']
                    eval(c_0_0).font.italic = eval(c_1_0).font.italic = epg_format_sbc['font_italic']
                    eval(c_0_0).font.bold = eval(c_1_0).font.bold = epg_format_sbc['font_bold'] 
                    eval(c_0_0).font.size = eval(c_1_0).font.size = epg_format_sbc['font_size'] 
                    eval(c_0_0).font.underline = eval(c_1_0).font.underline = epg_format_sbc['font_underline']
                    eval(c_0_0).alignment = eval(c_1_0).font.alignment = epg_format_sbc['allineamento']
                
    #SALVA PPTX OUTPUT
    vpdc.save(join(scriptdir, "03_DOC_NETWORK", dict_deliv_form['Ragione_Sociale_Cliente'] + ' - VPDC UCC.pptx'))
    


#COMPILA FOGLIO EXCEL
def config_compila_excel():
    #NOTE COMPILA FOGLIO EXCEL
    
    os.chdir(join(scriptdir, "Templates", "xlsx"))
    
    xlsx = "Template - vFW Edge vCloud.xlsx"

    xlsx = load_workbook(xlsx)
    
    xlsx_sheet = xlsx['vFirewall']
    
    
    #INSERIMENTO DATI
    
    #TITOLO (PRIMA RIGA)
    xlsx_sheet['B1'] = 'Firewall vCloud ' + dict_deliv_form['Ragione_Sociale_Cliente']
    
    #TITOLO TABELLA1
    xlsx_sheet['B3'] = 'Virtual firewall ' + dict_deliv_form['Ragione_Sociale_Cliente']
    
    #FW PORTA1
    fw_porta1= {'ip_address1': 'D5', 'default_gateway1':['E5','E18'],'vlan_id1': 'F5', 'port_profile1':'G5'}
    
    xlsx_sheet[fw_porta1['ip_address1']] = dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_NSXEDGE'] + " " + dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_NETMASK']
    xlsx_sheet[fw_porta1['default_gateway1'][0]] = dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_HSRP']
    xlsx_sheet[fw_porta1['vlan_id1']] = "vlan-" + dict_deliv_form['Vlan_internet'] + "- no tagged"
    xlsx_sheet[fw_porta1['port_profile1']] = dict_deliv_form['epg_internet_xlsx']
    
    #FW PORTA2
    fw_porta2= {'ip_address2': 'D6', 'default_gateway2':['E6','E19'],'vlan_id2': 'F6', 'port_profile2':'G6'}
    
    xlsx_sheet[fw_porta2['ip_address2']] = dict_deliv_form['Subnet_Interconnessione_vFW_UCC'] + " " + dict_deliv_form['Subnet_Interconnessione_NETMASK']
    xlsx_sheet[fw_porta2['default_gateway2'][0]] = dict_deliv_form['Subnet_Interconnessione_HSRP']
    xlsx_sheet[fw_porta2['vlan_id2']] = "vlan-" + dict_deliv_form['Vlan_Subnet_Interconnessione']
    xlsx_sheet[fw_porta2['port_profile2']] = dict_deliv_form['epg_vpn_mpls_xlsx']
    
    #FW ROUTING
    xlsx_sheet[fw_porta1['default_gateway1'][1]] = dict_deliv_form['Subnet_IP_pubblica_staccata_su_SOM_HSRP']
    xlsx_sheet[fw_porta2['default_gateway2'][1]] = dict_deliv_form['Subnet_Interconnessione_HSRP']
    
    # Save the file
    xlsx.save(join(scriptdir, "03_DOC_NETWORK", dict_deliv_form['Ragione_Sociale_Cliente'] + " - vFW Edge vCloud.xlsx")) 
    
    
#TODO? USER INPUT
def user_input():
    #NOTE CHIEDI VLAN INTERNET, MPLS, ITCFONIA
    #NOTE CHIEDI 71/72 <SEQ> prefix RedStatic/RedConnected INTERNET (è lo stesso?) 
    pass
    

###MAIN####
def main():
    # deliv_form_to_dict()
    # checklist_dati()
    user_input()
    txt_template_to_j2()
    deliv_form_integrazione()
    jinja_config()
    compila_pptx()
    config_compila_excel()

main()








#####################APPOGGIO


# test = df[df['CAMPO'].isin(['Seconda_Subnet_IP_30_MPLS_Fonia'])]
# print(test)

#ESTRAZIONE CELLA COLONNA VALORI DA STRINGA COLONNA CAMPO
#print(df.query('CAMPO == "Ragione_Sociale_Cliente"')['VALORE'].item()) #METODO 1
#print(df.loc[df['CAMPO'] == 'Ragione_Sociale_Cliente','VALORE'].item()) #METODO 2
#print(df.loc[df['CAMPO'] == 'Ragione_Sociale_Cliente']['VALORE'].values[0]) #METODO 3
#df.loc[df['CAMPO'] == 30000, 'VALORE'] 

#CONVERTE DF IN DICT
#dict_deliv_form = dict(zip(df.CAMPO, df.VALORE))
#dict_autom = df.set_index('CAMPO')['VALORE'].to_dict()   #alternativa a zip
#print(dictauto['Prima_Subnet_IP_/30_(MPLS_Fonia)'])

# def main():

    # # SCELTA PARTE DELLO SCRIPT (PE, PRE-ATTIVAZIONE O POST-ATTIVAZIONE, SESSIONI MPManager)
    # done = 0

    # while done == 0:

    #     print("\n\n")
    #     print("+-----------------------------------+")
    #     print("+---- SCELTA SCRIPT DA ESEGUIRE ----+")
    #     print("+-----------------------------------+")

    #     print("ALL: tutti i moduli")
    #     print("PROGETTO: solo il PE")
    #     print(
    #         "PRE-ATTIVAZIONE: tutte le le email, moduli attivazione, collaudo, postattivazione"
    #     )
    #     print("POST-ATTIVAZIONE: modulo collaudo compilato")
    #     print("PUTTY: sessioni Multi Putty Manager\n")

    #     script = input(
    #         "ALL [all]? [...|...|]: "
    #     ).lower()

    #     if script == "all":
    #         done = 1
    #         all()

    #     elif script == "...":
    #         done = 1
    #         pe()


    #     # USERINPUT ERRATO
    #     else:
    #         done = 0


# #CARICA DELIVERY FORM IN dict 
# #NOTE SI POTREBBE ELIMINARE
# def deliv_form_to_dict():
    
#     #CREA DICT DI BASE CON DATI PRESENTI IN DF DELIVERY FORM
    
#     global dict_deliv_form
#     df_filter_list = []
    
#     #crea lista df_filter_list per filtrare i campi del file di configurazione dal df pandas
#     for diz in campi_list_of_dicts:
#         # print(diz['filename'])
#         # if 'IN_MODIFICA-Cliente_ACI-VPDC_FONIA_L3_vrf_voce' in diz['filename']: #o diz.values()
#             for k,v in diz.items():
#                 # print(k,": ", v)
#                 df_filter_list.append(k)
    
    
#     config_file_df = df
#     #crea df filtrato per configurazione 
#     # config_file_df = df[df['CAMPO'].isin(df_filter_list)]
#     config_file_df = config_file_df.set_index('CAMPO') #elimina indice numerico
#     # print(config_file_df)
    
#     #crea il dizionario dal df filtrato
#     for entry in config_file_df.columns:
#         # print(entry)
#         for tupla in config_file_df[entry].items(): #.keys(), .values, .items()
#             dict_deliv_form[tupla[0]] = tupla[1]
        
#     # pprint(dict_deliv_form, sort_dicts=False)
